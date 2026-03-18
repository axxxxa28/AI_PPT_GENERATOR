import os 
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv

class SlideGenerator:
    def __init__(self, output_dir="ouput/"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)
        
    def create_slide_deck(self, title,bullet_points):
        prs= Presentation()

        #title slide
        slide_layout = prs.slide_layouts[0]  # Using the first layout (title
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        #content slide 

        for i in range(0, len(bullet_points), 5):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]
            title_shape.text = f"{title} (part {i//5 + 1})"
            for bullet in bullet_points[i:i+5]:
                p = content_shape.text_frame.add_paragraph()
                p.text = bullet
                
        # Save the presentation

        pptx_path = os.path.join(self.output_dir, "generated_presentation.pptx")
        prs.save(pptx_path)
        print(f"\n presentation saved at : {pptx_path}")

if __name__ == "__main__":
    #ex input 

    sample_title = "Sample Presentation"
    sample_points = [
        "Introduction to the topic",
        "Key concepts and definitions",
        "Detailed analysis of the subject",
        "Case studies and examples",
        "Conclusion and future directions",
        "Q&A session",
        "References and further reading",
        "Acknowledgments",
        "Contact information",
        "Thank you for your attention"
    ]

    slide_generator = SlideGenerator()
    slide_generator.create_slide_deck(sample_title, sample_points)
